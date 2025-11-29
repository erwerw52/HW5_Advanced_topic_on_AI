from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from typing import Dict, List
import os


def hex_to_rgb(hex_color: str):
    hex_color = hex_color.lstrip('#')
    lv = len(hex_color)
    return tuple(int(hex_color[i:i+lv//3], 16) for i in range(0, lv, lv//3))


def create_pptx(slides: List[Dict], template: Dict, output_path: str):
    prs = Presentation()
    # define title slide layout
    title_slide_layout = prs.slide_layouts[0]

    # Title slide
    first = slides[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = first.get('title', '')
    subtitle.text = first.get('subtitle', '')

    # apply style to title
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = template.get('font')
                run.font.size = Pt(template.get('title_size'))
                r, g, b = hex_to_rgb(template.get('primary_color'))
                run.font.color.rgb = RGBColor(r, g, b)

    # Content slides
    for s in slides[1:]:
        if s['type'] == 'content':
            slide_layout = prs.slide_layouts[1]  # title & content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = s.get('heading', '')
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            bullets = s.get('bullets', [])
            if bullets:
                for i, b in enumerate(bullets):
                    p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
                    p.text = b
                    p.level = 0
                    for run in p.runs:
                        run.font.name = template.get('font')
                        run.font.size = Pt(template.get('content_size'))

        elif s['type'] == 'ending':
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = s.get('title', '感謝聆聽，敬請指教')

    # Save
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    prs.save(output_path)
    return output_path
