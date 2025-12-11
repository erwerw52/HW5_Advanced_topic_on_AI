from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import io
import os

def analyze_input_slide(slide, slide_index, total_slides):
    """分析輸入投影片的內容結構"""
    info = {
        'slide_index': slide_index,
        'is_first': slide_index == 0,
        'is_last': slide_index == total_slides - 1,
        'layout_name': slide.slide_layout.name,
        'has_title': False,
        'title_text': '',
        'placeholder_count': 0,
        'text_shapes': [],
        'image_shapes': [],
        'other_shapes': []
    }
    
    # 檢查標題
    try:
        if slide.shapes.title and slide.shapes.title.text:
            info['has_title'] = True
            info['title_text'] = slide.shapes.title.text
    except:
        pass
    
    # 分析所有形狀
    for shape in slide.shapes:
        if shape.is_placeholder:
            info['placeholder_count'] += 1
        
        if shape.has_text_frame and shape.text.strip():
            info['text_shapes'].append({
                'text': shape.text,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'is_title': (shape == slide.shapes.title if hasattr(slide.shapes, 'title') else False)
            })
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            info['image_shapes'].append({
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            })
    
    return info

def select_template_slide(slide_info, template_slides):
    """根據投影片內容選擇最適合的模板投影片"""
    has_title = slide_info['has_title']
    text_count = len([s for s in slide_info['text_shapes'] if not s['is_title']])
    image_count = len(slide_info['image_shapes'])
    is_first = slide_info['is_first']
    is_last = slide_info['is_last']
    total_template_slides = len(template_slides)
    
    # 最後一張：使用模板的最後一張
    if is_last:
        return total_template_slides - 1
    
    # 第一張投影片使用模板的第一張（通常是標題頁）
    if is_first:
        return 0
    
    # 根據內容數量循環使用模板投影片
    # 模板有14張投影片（索引 0-13）
    # 我們從第2張開始循環使用（跳過第1張標題頁和最後1張結束頁）
    available_slides = list(range(1, total_template_slides - 1))  # [1, 2, 3, ..., 12]
    
    # 根據投影片索引循環選擇
    slide_index = slide_info['slide_index']
    # 去掉第一張後的索引
    adjusted_index = slide_index - 1
    
    # 循環使用可用的模板投影片
    selected_index = available_slides[adjusted_index % len(available_slides)]
    
    return selected_index

def get_template_placeholders(slide):
    """獲取模板投影片中的佔位符"""
    placeholders = {
        'title': [],
        'content': []
    }
    
    for shape in slide.placeholders:
        try:
            ptype = shape.placeholder_format.type
            
            # TITLE (1) 或 CENTER_TITLE (3) 視為標題
            if ptype == PP_PLACEHOLDER.TITLE or ptype == PP_PLACEHOLDER.CENTER_TITLE:
                placeholders['title'].append({
                    'shape': shape,
                    'top': shape.top,
                    'left': shape.left,
                    'idx': shape.placeholder_format.idx
                })
            # BODY (2), OBJECT (7), SUBTITLE (4) 視為內容
            elif ptype in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE]:
                placeholders['content'].append({
                    'shape': shape,
                    'top': shape.top,
                    'left': shape.left,
                    'idx': shape.placeholder_format.idx
                })
        except:
            pass
    
    # 按位置排序（先上後下，先左後右）
    placeholders['title'].sort(key=lambda x: (x['top'], x['left']))
    placeholders['content'].sort(key=lambda x: (x['top'], x['left']))
    
    return placeholders

def copy_text_to_shape(source_text, target_shape, is_title=False):
    """將文字複製到目標形狀並調整格式"""
    if not target_shape.has_text_frame:
        return False
    
    target_shape.text_frame.clear()
    target_shape.text = source_text
    
    # 調整文字格式
    for paragraph in target_shape.text_frame.paragraphs:
        if is_title:
            # 標題置中對齊
            paragraph.alignment = PP_ALIGN.CENTER
        
        for run in paragraph.runs:
            if not is_title:
                # 內容文字縮小
                if run.font.size and run.font.size > Pt(14):
                    run.font.size = Pt(14)
                elif not run.font.size:
                    run.font.size = Pt(14)
    
    return True

def copy_all_shapes_from_template(template_slide, new_slide):
    """從模板投影片複製所有形狀（包括佔位符）到新投影片"""
    copied_shapes = []
    for shape in template_slide.shapes:
        try:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            copied_shapes.append(shape.name)
        except Exception as e:
            print(f"    ⚠ 無法複製形狀 {shape.name}: {e}")
    return copied_shapes

def copy_images_from_input(input_slide, new_slide):
    """從輸入投影片複製圖片到新投影片"""
    images_copied = 0
    for shape in input_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                # 獲取圖片資訊
                image = shape.image
                image_bytes = image.blob
                
                # 在新投影片中添加圖片
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # 使用圖片的二進制數據創建新圖片
                pic = new_slide.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    left, top, width, height
                )
                images_copied += 1
            except Exception as e:
                print(f"    ⚠ 無法複製圖片: {e}")
    
    return images_copied

def create_from_template(input_path, template_path, output_path):
    """讀取輸入PPT和模板PPT，將內容套用到模板生成新PPT"""
    print(f"\n=== 開始處理 ===")
    print(f"輸入檔案: {input_path}")
    print(f"模板檔案: {template_path}")
    
    # 1. 讀取輸入PPT
    input_prs = Presentation(input_path)
    print(f"\n讀取輸入PPT: 共 {len(input_prs.slides)} 張投影片")
    
    # 2. 讀取模板PPT
    template_prs = Presentation(template_path)
    print(f"讀取模板PPT: 共 {len(template_prs.slide_layouts)} 種佈局")
    
    # 3. 創建基於模板的新簡報
    output_prs = Presentation(template_path)
    
    # 移除模板中的現有投影片
    xml_slides = output_prs.slides._sldIdLst
    slides = list(xml_slides)
    for s in slides:
        xml_slides.remove(s)
    
    print(f"\n開始轉換投影片...")
    
    # 4. 逐張處理輸入投影片
    for i, slide in enumerate(input_prs.slides):
        print(f"\n處理投影片 {i+1}/{len(input_prs.slides)}")
        
        # 分析投影片內容
        slide_info = analyze_input_slide(slide, i, len(input_prs.slides))
        print(f"  原始佈局: {slide_info['layout_name']}")
        print(f"  標題: {slide_info['title_text'][:50]}..." if slide_info['title_text'] else "  標題: 無")
        print(f"  文字區塊: {len(slide_info['text_shapes'])}")
        print(f"  圖片: {len(slide_info['image_shapes'])}")
        
        # 最後一頁：保留原始模板的最後一頁
        if slide_info['is_last']:
            print(f"  這是最後一頁，使用原始模板的最後一張投影片")
            # 使用模板的最後一張投影片（如果有的話）
            if len(template_prs.slides) > 0:
                # 複製模板的最後一張投影片
                template_last_slide = template_prs.slides[-1]
                slide_layout = template_last_slide.slide_layout
                new_slide = output_prs.slides.add_slide(slide_layout)
                
                # 複製模板投影片的所有形狀（包括佔位符和裝飾）
                copy_all_shapes_from_template(template_last_slide, new_slide)
                
                # 只複製標題（如果有）
                if slide_info['has_title']:
                    try:
                        if new_slide.shapes.title:
                            new_slide.shapes.title.text = slide_info['title_text']
                            copy_text_to_shape(slide_info['title_text'], new_slide.shapes.title, is_title=True)
                            print(f"  >> 已複製標題到模板最後一頁")
                    except:
                        pass
                
                # 複製輸入投影片的圖片
                images_copied = copy_images_from_input(slide, new_slide)
                if images_copied > 0:
                    print(f"  >> 已複製 {images_copied} 張圖片")
            continue
        
        # 選擇合適的模板投影片
        template_slide_index = select_template_slide(slide_info, template_prs.slides)
        template_slide = template_prs.slides[template_slide_index]
        
        print(f"  使用模板投影片: 第{template_slide_index + 1}張 - {template_slide.slide_layout.name}")
        
        # 使用該模板投影片的佈局創建新投影片
        slide_layout = template_slide.slide_layout
        new_slide = output_prs.slides.add_slide(slide_layout)
        
        # 複製模板投影片的所有形狀（包括佔位符和裝飾）
        copied_shapes = copy_all_shapes_from_template(template_slide, new_slide)
        print(f"  >> 已複製 {len(copied_shapes)} 個形狀")
        
        # 準備要填入的內容
        input_title = slide_info['title_text'] if slide_info['has_title'] else ''
        input_contents = [s['text'] for s in slide_info['text_shapes'] if not s['is_title']]
        
        print(f"  輸入標題: {'有' if input_title else '無'}")
        print(f"  輸入內容: {len(input_contents)} 個")
        
        # 獲取新投影片的文字形狀並替換內容
        text_shapes = [s for s in new_slide.shapes if s.has_text_frame]
        title_shapes = []
        content_shapes = []
        
        for shape in text_shapes:
            # 檢查是否有文字內容
            has_text = shape.text.strip() != ''
            
            if shape.is_placeholder:
                try:
                    ptype = shape.placeholder_format.type
                    if ptype == PP_PLACEHOLDER.TITLE or ptype == PP_PLACEHOLDER.CENTER_TITLE:
                        title_shapes.append(shape)
                    else:
                        content_shapes.append(shape)
                except:
                    content_shapes.append(shape)
            else:
                # 非佔位符的文字框
                # 根據文字內容判斷是標題還是內容
                if has_text:
                    # 如果文字較短且位置靠上，視為標題候選
                    if len(shape.text) < 100 and shape.top < Inches(2):
                        title_shapes.append(shape)
                    else:
                        content_shapes.append(shape)
                else:
                    content_shapes.append(shape)
        
        # 按位置排序
        title_shapes.sort(key=lambda x: (x.top, x.left))
        content_shapes.sort(key=lambda x: (x.top, x.left))
        
        # 替換標題文字
        title_replaced = 0
        used_title_shapes = []
        if input_title:
            if len(title_shapes) > 0:
                # 優先使用第一個標題形狀
                title_shapes[0].text = input_title
                copy_text_to_shape(input_title, title_shapes[0], is_title=True)
                used_title_shapes.append(title_shapes[0])
                title_replaced += 1
                print(f"  >> 已替換標題文字")
        
        # 替換內容文字
        content_replaced = 0
        used_content_shapes = []
        for i, content_text in enumerate(input_contents):
            if i < len(content_shapes):
                content_shapes[i].text = content_text
                copy_text_to_shape(content_text, content_shapes[i], is_title=False)
                used_content_shapes.append(content_shapes[i])
                content_replaced += 1
        
        if content_replaced > 0:
            print(f"  >> 已替換 {content_replaced} 個內容文字")
        elif len(input_contents) > 0 and len(content_shapes) == 0:
            print(f"  !! 警告: 有 {len(input_contents)} 個內容但模板沒有內容區域")
        
        # 移除未使用的佔位符和文字框
        removed_count = 0
        shapes_to_remove = []
        
        for shape in title_shapes:
            if shape not in used_title_shapes:
                shapes_to_remove.append(shape)
        
        for shape in content_shapes:
            if shape not in used_content_shapes:
                shapes_to_remove.append(shape)
        
        # 執行移除
        for shape in shapes_to_remove:
            try:
                sp = shape.element
                sp.getparent().remove(sp)
                removed_count += 1
            except Exception as e:
                pass
        
        if removed_count > 0:
            print(f"  >> 已移除 {removed_count} 個未使用的佔位符")
        
        # 複製輸入投影片的圖片
        images_copied = copy_images_from_input(slide, new_slide)
        if images_copied > 0:
            print(f"  >> 已複製 {images_copied} 張圖片")
    
    # 5. 儲存輸出檔案
    output_prs.save(output_path)
    print(f"\n=== 完成 ===")
    print(f"輸出檔案: {output_path}")

